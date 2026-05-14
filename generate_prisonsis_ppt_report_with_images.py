"""
生成《PrisonSIS_产品介绍_汇报版.pptx》：领导汇报口径 + 嵌入 docs/ppt-screenshots/汇报版 下 PNG。

前置：在 frontend 目录执行 `npm run ppt:screenshots`（需本机已 `pip install python-pptx`，且已 `npx playwright install chromium`）。
"""

from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parent
IMG_DIR = REPO_ROOT / "docs" / "ppt-screenshots" / "汇报版"
OUTPUT = REPO_ROOT / "PrisonSIS_产品介绍_汇报版.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)


def add_image_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    caption: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.65))
    tfp = title_box.text_frame
    tfp.clear()
    p0 = tfp.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(28)
    p0.font.bold = True

    top = Inches(1.05)
    left = Inches(0.45)
    pic_w = Inches(12.4)
    pic_h = Inches(5.85)
    slide.shapes.add_picture(str(image_path), left, top, width=pic_w, height=pic_h)

    if caption:
        cap = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.45))
        ctf = cap.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(12)
        cp.font.italic = True


def add_quad_screenshot_slide(
    prs: Presentation,
    title: str,
    paths: list[tuple[str, Path]],
) -> None:
    """paths: (label, path) 最多 4 个，缺失则跳过。"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.65))
    tf = tb.text_frame
    tf.clear()
    pr = tf.paragraphs[0]
    pr.text = title
    pr.font.size = Pt(28)
    pr.font.bold = True

    existing = [(lab, p) for lab, p in paths if p.exists()]
    if not existing:
        body = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(12), Inches(2))
        btf = body.text_frame
        btf.clear()
        bp = btf.paragraphs[0]
        bp.text = "未找到截图文件。请在 frontend 目录执行：npm run ppt:screenshots"
        bp.font.size = Pt(18)
        return

    # 2x2 网格（英寸坐标，避免对 Length 做非常规运算）
    gap = Inches(0.12)
    margin_l = Inches(0.45)
    margin_t = Inches(1.05)
    cell_w = Inches(6.15)
    cell_h = Inches(2.85)
    positions = [
        (margin_l, margin_t),
        (margin_l + cell_w + gap, margin_t),
        (margin_l, margin_t + cell_h + gap),
        (margin_l + cell_w + gap, margin_t + cell_h + gap),
    ]

    for idx, (label, img_path) in enumerate(existing[:4]):
        lx, ty = positions[idx]
        slide.shapes.add_picture(str(img_path), lx, ty, width=cell_w, height=cell_h)
        cap = slide.shapes.add_textbox(lx, ty + cell_h + Inches(0.02), cell_w, Inches(0.22))
        ctf = cap.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp.text = label
        cp.font.size = Pt(11)


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    today = date.today().strftime("%Y-%m-%d")

    add_title_slide(
        prs,
        "监狱审讯笔录系统（PrisonSIS）",
        f"产品介绍（汇报版）\n阶段：开发中 · 版本 v1.0 · {today}\n界面截图为示例数据，仅供演示",
    )

    add_bullet_slide(
        prs,
        "1. 背景与痛点",
        [
            "笔录格式与填写习惯不统一，复核成本高、质量波动大",
            "审批协同依赖口头沟通，进度不透明、难以及时督办",
            "审计与追责依赖零散记录，检索与导出效率偏低",
        ],
    )

    add_bullet_slide(
        prs,
        "2. 建设目标与价值",
        [
            "目标：形成「录入 — 审批 — 归档 — 审计」闭环，支撑日常办案与内控",
            "价值：规范留痕、提升协同效率、降低审计与复盘成本",
            "形态：桌面端优先，本地数据与权限可控，便于内网部署",
        ],
    )

    add_bullet_slide(
        prs,
        "3. 功能版图（全模块）",
        [
            "核心业务：首页、罪犯信息、笔录制作、审批中心、案件管理、档案管理",
            "资源配置：统计分析、模板管理、文档导出、用户管理、数据备份",
            "系统管理：日志审计；并支持全局搜索跨模块快速定位（能力随版本持续完善）",
            "规划：罪犯信息后续支持内网权威库检索与批量导入，减少重复录入与差错",
        ],
    )

    add_bullet_slide(
        prs,
        "4. 路线图（占位，可按实际计划调整）",
        [
            "当前：核心流程与权限框架开发中，持续联调与体验优化",
            "下一里程碑：试点环境部署与真实数据演练（仍须脱敏与合规评审）",
            "后续：统计分析深化、与其他内网系统对接（含罪犯信息检索导入）",
            "推广：在稳定运行与培训体系就绪后，分阶段扩大使用范围",
        ],
    )

    add_quad_screenshot_slide(
        prs,
        "5. 核心界面（拼图）",
        [
            ("笔录制作", IMG_DIR / "records.png"),
            ("审批中心", IMG_DIR / "approvals.png"),
            ("日志审计", IMG_DIR / "logs.png"),
            ("全局搜索", IMG_DIR / "global-search.png"),
        ],
    )

    add_bullet_slide(
        prs,
        "6. 技术架构（概要）",
        [
            "客户端：桌面应用，统一入口与操作体验，适配日常办公场景",
            "数据层：本地结构化存储，便于离线可用与快速检索",
            "工程实现：现代前端 + 本地高性能运行时；具体技术栈见研发材料",
        ],
    )

    add_bullet_slide(
        prs,
        "7. 安全与权限",
        [
            "分级权限：管理员、审计、审批、经办等角色分权分域",
            "关键能力受控：日志清理、备份恢复、用户管理等仅授权角色可操作",
            "审计留痕：关键动作可追溯（用户、时间、对象与摘要信息）",
        ],
    )

    home_png = IMG_DIR / "home.png"
    if home_png.exists():
        add_image_slide(
            prs,
            "8. 整体界面一览",
            home_png,
            caption="侧栏为各业务模块入口；顶栏含主题切换与全局搜索（示例数据）",
        )
    else:
        add_bullet_slide(
            prs,
            "8. 整体界面一览",
            ["请先执行 npm run ppt:screenshots 生成 home.png 后重新导出本页。"],
        )

    add_bullet_slide(
        prs,
        "9. 合规说明与总结陈述",
        [
            "合规：材料与截图为示例数据，不含真实单位、人员与案件；对外演示须按规定脱敏与审批",
            "总结：聚焦笔录业务闭环，当前开发迭代中，主线能力可演示、可持续打磨",
            "建议：明确试点范围、数据接入（含罪犯信息内网检索导入）与验收指标；恳请支持资源与推广节奏",
        ],
    )

    prs.save(OUTPUT)
    print(f"已写入: {OUTPUT}")


if __name__ == "__main__":
    main()
