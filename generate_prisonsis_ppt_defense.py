from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(22)


def add_qa_slide(prs, title, qa_pairs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12.0), Inches(5.7))
    tf = box.text_frame
    tf.clear()

    for idx, (q, a) in enumerate(qa_pairs):
        qp = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        qp.text = f"Q{idx + 1}: {q}"
        qp.font.bold = True
        qp.font.size = Pt(20)

        ap = tf.add_paragraph()
        ap.text = f"A: {a}"
        ap.level = 1
        ap.font.size = Pt(18)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "PrisonSIS 项目答辩版",
        "监狱审讯笔录工具（含Q&A备页）",
    )

    add_bullet_slide(prs, "1. 课题背景", [
        "监狱审讯笔录流程长期依赖人工与线下沟通",
        "存在模板不统一、审批效率低、审计追溯难等问题",
        "需要一套本地化、安全可控的数字化解决方案",
    ])

    add_bullet_slide(prs, "2. 项目目标", [
        "规范化：统一模板与录入标准",
        "高效率：缩短审批与协同链路",
        "可追溯：全链路日志审计与导出",
    ])

    add_bullet_slide(prs, "3. 系统方案", [
        "前端：React + TypeScript + 玻璃主题交互",
        "桌面容器：Tauri",
        "后端：Rust + SQLite（本地数据存储）",
    ])

    add_bullet_slide(prs, "4. 核心功能", [
        "笔录制作：模板套用、引导式录入、富文本大纲定位",
        "审批中心：待办处理、通过/驳回、状态可视化",
        "日志审计：筛选检索、导出留痕、责任可追踪",
    ])

    add_bullet_slide(prs, "5. 关键优化成果", [
        "实现全局搜索面板：跨 6 模块聚合命中",
        "统一深浅主题，修复多处可读性与层级问题",
        "优化编辑/查看态大纲跳转与高亮一致性",
    ])

    add_bullet_slide(prs, "6. 项目价值", [
        "业务价值：提升录入质量与审批时效",
        "管理价值：降低审计检索与追责成本",
        "技术价值：形成可扩展的桌面端业务框架",
    ])

    add_bullet_slide(prs, "7. 风险与对策", [
        "风险：权限误操作、数据恢复风险、UI复杂性上升",
        "对策：角色分级、关键操作确认、审计留痕闭环",
        "保障：持续回归验证与逐步发布策略",
    ])

    add_bullet_slide(prs, "8. 下一步计划", [
        "完善统计分析与报表能力",
        "增加自动化测试与性能基线",
        "推进试点推广与用户反馈闭环",
    ])

    add_qa_slide(prs, "Q&A 备页 1：业务类问题", [
        ("这套系统最直接提升了什么？", "笔录规范性、审批效率、审计追溯速度。"),
        ("为什么要做全局搜索？", "减少跨页面检索成本，快速定位业务对象。"),
        ("是否支持线下部署？", "支持，Tauri + SQLite 适合本地化部署场景。"),
    ])

    add_qa_slide(prs, "Q&A 备页 2：技术类问题", [
        ("为什么选 Rust + Tauri？", "兼顾性能、安全与桌面端本地能力。"),
        ("如何保证权限安全？", "按角色分级授权，敏感操作留痕可审计。"),
        ("后续扩展怎么做？", "可沿现有模块化结构增加报表、检索、接口层能力。"),
    ])

    add_bullet_slide(prs, "致谢", [
        "感谢各位老师/评审指导",
        "欢迎提问与交流",
    ])

    output_path = "PrisonSIS_答辩版_含QA备页.pptx"
    prs.save(output_path)
    print(output_path)


if __name__ == "__main__":
    main()
