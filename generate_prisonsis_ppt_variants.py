from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(22)


def new_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def build_exec_brief():
    prs = new_presentation()
    add_title_slide(
        prs,
        "PrisonSIS 领导汇报精简版",
        "监狱审讯笔录工具（8页）",
    )
    add_bullet_slide(prs, "1. 项目价值", [
        "规范化笔录全流程，减少人为差异",
        "审批与审计数字化闭环，提升管理效率",
        "本地化部署，安全可控、便于推广",
    ])
    add_bullet_slide(prs, "2. 痛点与风险", [
        "模板不统一，录入质量不稳定",
        "审批链断点多，状态不可视",
        "日志检索与追溯成本高",
    ])
    add_bullet_slide(prs, "3. 方案与闭环", [
        "录入：模板 + 引导式 + 富文本大纲定位",
        "审批：提交、待办、通过/驳回全流程可视",
        "审计：日志筛选、导出、留痕追责",
    ])
    add_bullet_slide(prs, "4. 核心能力", [
        "笔录制作、审批中心、案件管理、档案管理",
        "日志审计、用户权限、导出备份",
        "全局搜索跨模块联动跳转",
    ])
    add_bullet_slide(prs, "5. 阶段成果", [
        "全局搜索面板已上线（跨6模块）",
        "深浅主题适配并修复关键可读性问题",
        "编辑态/查看态大纲定位体验统一",
    ])
    add_bullet_slide(prs, "6. 量化收益（建议填真实数据）", [
        "笔录录入效率提升：xx%",
        "审批处理时效提升：xx%",
        "审计定位时间下降：xx%",
    ])
    add_bullet_slide(prs, "7. 风险与对策", [
        "风险：权限误用、数据恢复误操作、交互复杂度",
        "对策：角色控制 + 审计留痕 + 分级确认",
        "持续迭代：体验细化与性能优化并行推进",
    ])
    add_bullet_slide(prs, "8. 下一步与资源诉求", [
        "完善统计分析与报表能力",
        "推进试点落地，收集业务反馈",
        "建议保障迭代节奏：产品/前端/后端协同投入",
    ])
    return prs


def build_tech_review():
    prs = new_presentation()
    add_title_slide(
        prs,
        "PrisonSIS 技术评审版",
        "架构与实现细节（10页）",
    )
    add_bullet_slide(prs, "1. 需求边界", [
        "目标：录入-审批-归档-审计全链路数字化",
        "范围：桌面端业务系统，角色分级权限",
        "非目标：移动端原生、跨机构联邦部署",
    ])
    add_bullet_slide(prs, "2. 业务流程", [
        "笔录创建/编辑 -> 提交审批 -> 审批结果",
        "已归档/在押切换管理",
        "日志审计检索与导出",
    ])
    add_bullet_slide(prs, "3. 技术栈", [
        "前端：React + TypeScript + 玻璃主题系统",
        "容器：Tauri（桌面端）",
        "后端：Rust + SQLite（本地数据）",
    ])
    add_bullet_slide(prs, "4. 分层架构", [
        "UI 层：页面与组件（pages/components）",
        "服务层：API bridge 与聚合逻辑（api/lib）",
        "后端命令层：Tauri invoke -> Rust command -> DB",
    ])
    add_bullet_slide(prs, "5. 数据模型", [
        "核心实体：Record / Criminal / Case / Template / AuditLog",
        "类型定义集中在 frontend/src/api/types.ts",
        "分页接口统一返回 [rows, total]",
    ])
    add_bullet_slide(prs, "6. 全局搜索实现", [
        "入口：GlassHeader 顶栏搜索",
        "聚合：runGlobalSearch 并发查询 6 模块",
        "展示：GlobalSearchPanel 分组结果 + 组内滚动",
    ])
    add_bullet_slide(prs, "7. 页面联动机制", [
        "App 层统一状态：open/loading/error/query/results",
        "事件总线：prisonsis:apply-search",
        "各页监听并回填 searchInput/appliedSearch",
    ])
    add_bullet_slide(prs, "8. 主题与可读性", [
        "深浅主题变量化（index.css）",
        "状态栏、侧栏收起态、搜索面板等定点优化",
        "减少硬编码 rgba，提升一致性与维护性",
    ])
    add_bullet_slide(prs, "9. 质量与验证", [
        "构建校验：npm run build",
        "关键交互手测：跳转、高亮、滚动、搜索联动",
        "逐步收敛 UI 问题并保持最小改动策略",
    ])
    add_bullet_slide(prs, "10. 技术债与演进", [
        "全局搜索可升级为后端统一检索接口",
        "继续拆分样式 token，降低组件内联样式占比",
        "逐步补齐自动化测试与性能指标",
    ])
    return prs


def main():
    exec_path = "PrisonSIS_领导汇报精简版_8页.pptx"
    tech_path = "PrisonSIS_技术评审版_10页.pptx"
    build_exec_brief().save(exec_path)
    build_tech_review().save(tech_path)
    print(exec_path)
    print(tech_path)


if __name__ == "__main__":
    main()
