from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(22)


def add_two_col_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.9), Inches(5.0))
    right_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.6), Inches(5.9), Inches(5.0))

    left_tf = left_box.text_frame
    left_tf.clear()
    p = left_tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(26)
    p.font.bold = True
    for item in left_items:
        pp = left_tf.add_paragraph()
        pp.text = f"• {item}"
        pp.font.size = Pt(20)

    right_tf = right_box.text_frame
    right_tf.clear()
    p = right_tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(26)
    p.font.bold = True
    for item in right_items:
        pp = right_tf.add_paragraph()
        pp.text = f"• {item}"
        pp.font.size = Pt(20)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "监狱审讯笔录工具（PrisonSIS）",
        "产品介绍与阶段成果汇报\n版本：v1.0",
    )

    add_bullet_slide(
        prs,
        "1. 项目背景与痛点",
        [
            "笔录模板不统一，人工录入标准不一致",
            "审批链路依赖线下沟通，进度不可视",
            "日志追溯成本高，审计效率低",
        ],
    )

    add_bullet_slide(
        prs,
        "2. 产品定位与目标",
        [
            "定位：监狱场景下的笔录全流程数字化平台",
            "目标：规范录入、提升审批效率、强化审计追溯",
            "结果：形成“录入-审批-归档-审计”闭环",
        ],
    )

    add_bullet_slide(
        prs,
        "3. 核心能力总览",
        [
            "笔录制作：模板套用、引导式录入、富文本编辑",
            "审批中心：待办处理、通过/驳回、状态可视化",
            "管理能力：案件、档案、用户、日志、导出备份",
        ],
    )

    add_two_col_slide(
        prs,
        "4. 关键流程演示",
        "笔录制作",
        [
            "新建笔录并套用模板",
            "大纲点击定位 + 段落高亮",
            "编辑态与查看态体验统一",
        ],
        "审批与审计",
        [
            "提交审批并查看处理状态",
            "日志按用户/动作/模块筛选",
            "支持导出与追责留痕",
        ],
    )

    add_bullet_slide(
        prs,
        "5. 本阶段新增亮点",
        [
            "全局搜索面板：跨模块聚合命中并一键跳转",
            "深浅主题统一：白雾浅色与深色毛玻璃适配",
            "多处可读性修复：状态栏、侧栏收起、弹层滚动",
        ],
    )

    add_bullet_slide(
        prs,
        "6. 全局搜索能力（已上线）",
        [
            "搜索范围：笔录、罪犯、案件、审批、日志、模板",
            "展示方式：分组结果面板 + 组内滚动，避免穿模",
            "联动行为：点击结果后自动跳页并带入搜索词",
        ],
    )

    add_bullet_slide(
        prs,
        "7. 技术架构",
        [
            "前端：React + TypeScript + Glass UI",
            "桌面容器：Tauri",
            "后端：Rust + SQLite（本地高性能、易部署）",
        ],
    )

    add_bullet_slide(
        prs,
        "8. 安全与权限",
        [
            "角色分级：Admin / Auditor / Approver / User",
            "关键操作权限控制：清空日志、备份恢复、用户管理",
            "审计留痕：用户、动作、对象、时间、详情全链路记录",
        ],
    )

    add_bullet_slide(
        prs,
        "9. 下一步规划",
        [
            "完善统计分析与报表能力（业务指标可视化）",
            "持续优化交互细节与浅色主题一致性",
            "推进稳定性与性能优化，支持更大规模数据",
        ],
    )

    add_bullet_slide(
        prs,
        "10. 总结",
        [
            "PrisonSIS 已形成核心业务闭环并可实际使用",
            "当前重点从“可用”迈向“好用、稳用、可推广”",
            "建议进入试点推广与持续迭代阶段",
        ],
    )

    output_path = "PrisonSIS_产品介绍_汇报版.pptx"
    prs.save(output_path)
    print(output_path)


if __name__ == "__main__":
    main()
